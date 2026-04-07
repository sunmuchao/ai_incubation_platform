"""
文档加载器 - 支持 PDF, Word, Excel, PPT 文件解析
"""
import os
from typing import List, Dict, Any
from pathlib import Path

from .base import UnstructuredDataConnector, UnstructuredConfig, DocumentChunk

try:
    from utils.logger import logger
except ImportError:
    import logging
    logger = logging.getLogger(__name__)


class DocumentLoader(UnstructuredDataConnector):
    """文档加载器 - 支持 PDF, Word, Excel, PPT"""

    def __init__(self, config: UnstructuredConfig):
        super().__init__(config)
        self._document_text = ""
        self._metadata = {}

    async def connect(self) -> None:
        """验证文件是否存在"""
        if self.config.source_path:
            if not os.path.exists(self.config.source_path):
                raise FileNotFoundError(f"Document not found: {self.config.source_path}")
        self._connected = True
        logger.info("Document loader connected", extra={"source": self.config.source_path})

    async def disconnect(self) -> None:
        """断开连接，清理缓存"""
        self._document_text = ""
        self._metadata = {}
        self._content_cache = []
        self._connected = False
        logger.info("Document loader disconnected")

    async def load_content(self) -> List[DocumentChunk]:
        """加载文档内容"""
        if not self._connected:
            await self.connect()

        source_path = self.config.source_path
        if not source_path:
            raise ValueError("source_path is required for document loading")

        ext = Path(source_path).suffix.lower()

        if ext == '.pdf':
            text, metadata = await self._load_pdf(source_path)
        elif ext in ['.docx', '.doc']:
            text, metadata = await self._load_word(source_path)
        elif ext in ['.xlsx', '.xls']:
            text, metadata = await self._load_excel(source_path)
        elif ext in ['.pptx', '.ppt']:
            text, metadata = await self._load_ppt(source_path)
        else:
            raise ValueError(f"Unsupported document type: {ext}")

        self._document_text = text
        self._metadata = metadata

        # 分割成片段
        chunks = self._split_text(text)
        result = []
        for i, chunk in enumerate(chunks):
            result.append(self._create_chunk(
                content=chunk,
                index=i,
                metadata={**metadata, "chunk_type": "document", "file_type": ext}
            ))

        self._content_cache = result
        logger.info("Document loaded", extra={"path": source_path, "chunks": len(result)})
        return result

    async def get_schema(self) -> Dict[str, Any]:
        """获取文档元数据"""
        return {
            "source_type": "document",
            "source_path": self.config.source_path,
            "chunk_size": self.config.chunk_size,
            "chunk_overlap": self.config.chunk_overlap,
            "metadata": self._metadata,
            "total_chunks": len(self._content_cache)
        }

    async def _load_pdf(self, path: str) -> tuple:
        """加载 PDF 文件"""
        try:
            import PyPDF2

            text_parts = []
            metadata = {}

            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                metadata = reader.metadata or {}

                for i, page in enumerate(reader.pages):
                    text_parts.append(f"[Page {i + 1}]\n{page.extract_text()}")

            return "\n\n".join(text_parts), {
                "file_type": "pdf",
                "page_count": len(reader.pages),
                "metadata": {k: str(v) for k, v in metadata.items()}
            }
        except ImportError:
            # 尝试使用 pdfplumber 作为备选
            try:
                import pdfplumber

                text_parts = []
                with pdfplumber.open(path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        text_parts.append(f"[Page {i + 1}]\n{text}")

                return "\n\n".join(text_parts), {
                    "file_type": "pdf",
                    "page_count": len(text_parts)
                }
            except ImportError:
                raise ImportError(
                    "PyPDF2 or pdfplumber is required. "
                    "Install with: pip install PyPDF2 pdfplumber"
                )

    async def _load_word(self, path: str) -> tuple:
        """加载 Word 文件"""
        try:
            from docx import Document

            doc = Document(path)
            text_parts = []

            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    text_parts.append(para.text)

            # 添加表格内容
            for i, table in enumerate(doc.tables):
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text_parts.append(" | ".join(row_text))

            return "\n\n".join(text_parts), {
                "file_type": "docx",
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables)
            }
        except ImportError:
            raise ImportError(
                "python-docx is required. Install with: pip install python-docx"
            )

    async def _load_excel(self, path: str) -> tuple:
        """加载 Excel 文件"""
        try:
            import pandas as pd

            # 读取所有 sheet
            excel_file = pd.ExcelFile(path)
            text_parts = []
            sheet_info = {}

            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_parts.append(f"[Sheet: {sheet_name}]")
                text_parts.append(df.to_string())
                text_parts.append("")

                sheet_info[sheet_name] = {
                    "rows": len(df),
                    "columns": list(df.columns)
                }

            return "\n".join(text_parts), {
                "file_type": "xlsx",
                "sheet_count": len(excel_file.sheet_names),
                "sheets": sheet_info
            }
        except ImportError:
            raise ImportError(
                "pandas and openpyxl are required. "
                "Install with: pip install pandas openpyxl"
            )

    async def _load_ppt(self, path: str) -> tuple:
        """加载 PPT 文件"""
        try:
            from pptx import Presentation

            prs = Presentation(path)
            text_parts = []

            for i, slide in enumerate(prs.slides):
                text_parts.append(f"[Slide {i + 1}]")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            return "\n\n".join(text_parts), {
                "file_type": "pptx",
                "slide_count": len(prs.slides)
            }
        except ImportError:
            raise ImportError(
                "python-pptx is required. Install with: pip install python-pptx"
            )
